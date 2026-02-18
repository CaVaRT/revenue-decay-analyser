from __future__ import annotations

import argparse
import json
from datetime import datetime
import math
from pathlib import Path
from urllib.parse import urlparse

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots

from src.common.config import ROOT_DIR


SEVERITY_ORDER = ["Critical", "High", "Medium", "Low", "Monitor"]
QUADRANT_ORDER = [
    "Both Down",
    "Clicks Down, Revenue Stable/Up",
    "Revenue Down, Clicks Stable/Up",
    "Both Stable/Up",
]


def require_columns(df: pd.DataFrame, required: list[str], *, file_label: str) -> None:
    missing = [column for column in required if column not in df.columns]
    if missing:
        joined = ", ".join(missing)
        raise ValueError(f"{file_label} is missing required columns: {joined}")


def write_plot_image(fig, output_path: Path) -> bool:
    try:
        fig.write_image(str(output_path))
        return True
    except Exception as exc:  # pragma: no cover - runtime chart engine dependent
        print(f"Warning: could not write chart to {output_path}: {exc}")
        return False


def shorten_page(page: str, max_len: int = 55) -> str:
    text = str(page)
    if len(text) <= max_len:
        return text
    return text[: max_len - 3] + "..."


def normalize_page_path(raw: str) -> str:
    text = str(raw or "").strip()
    if not text:
        return ""
    if "://" in text:
        parsed = urlparse(text)
        path = parsed.path or "/"
    else:
        parsed = urlparse(f"https://dummy{text if text.startswith('/') else '/' + text}")
        path = parsed.path or "/"
    normalized = path.rstrip("/")
    return normalized or "/"


def pick_first_column(columns: list[str], candidates: list[str]) -> str | None:
    lowered = {column.lower(): column for column in columns}
    for candidate in candidates:
        found = lowered.get(candidate.lower())
        if found:
            return found
    return None


def to_bool_series(series: pd.Series) -> pd.Series:
    if pd.api.types.is_bool_dtype(series):
        return series.fillna(False)
    normalized = series.astype(str).str.strip().str.lower()
    return normalized.isin({"true", "1", "yes", "y", "t"})


def apply_base_layout(
    fig: go.Figure,
    *,
    left_margin: int = 200,
    height: int | None = None,
    width: int | None = None,
) -> go.Figure:
    fig.update_layout(
        margin={"l": left_margin, "r": 40, "t": 80, "b": 60},
        font={"size": 11},
        uniformtext_minsize=8,
        uniformtext_mode="hide",
        bargap=0.2,
    )
    fig.update_xaxes(automargin=True)
    fig.update_yaxes(automargin=True)
    if height:
        fig.update_layout(height=height)
    if width:
        fig.update_layout(width=width)
    return fig


def normalized_index(series: pd.Series) -> pd.Series:
    vals = pd.to_numeric(series, errors="coerce").fillna(0.0)
    base_candidates = vals[vals > 0]
    base = float(base_candidates.iloc[0]) if not base_candidates.empty else 1.0
    return (vals / base) * 100.0


def build_top_overall(scores: pd.DataFrame, *, top_n: int) -> tuple[pd.DataFrame, go.Figure]:
    top = scores.sort_values("decay_score", ascending=False).head(top_n).copy()
    top["page_short"] = top["page"].map(lambda value: shorten_page(value, max_len=70))
    top = top.sort_values("decay_score", ascending=True)
    fig = px.bar(
        top,
        x="decay_score",
        y="page_short",
        color="severity",
        orientation="h",
        hover_data=["page", "page_type", "total_revenue", "total_clicks", "losing_revenue_and_clicks"],
        category_orders={"severity": SEVERITY_ORDER},
        title=f"Top {top_n} Risk Pages Overall",
        labels={"decay_score": "Decay Score", "page_short": "Page"},
        height=max(650, 30 * len(top)),
    )
    fig.update_layout(
        yaxis={"categoryorder": "total ascending"},
        legend_title="Severity",
    )
    fig.update_yaxes(tickfont={"size": 10})
    apply_base_layout(fig, left_margin=320, width=1500)
    return top, fig


def build_top_per_page_type(
    scores: pd.DataFrame,
    *,
    top_n_per_type: int,
    max_page_types_facet: int,
) -> tuple[pd.DataFrame, go.Figure]:
    ordered = scores.sort_values("decay_score", ascending=False).copy()
    top_per_type = ordered.groupby("page_type", as_index=False, group_keys=False).head(top_n_per_type).copy()
    top_per_type["page_short"] = top_per_type["page"].map(lambda value: shorten_page(value, max_len=52))
    if top_per_type.empty:
        return top_per_type, go.Figure()

    selected_types = (
        top_per_type.groupby("page_type", as_index=False)["decay_score"]
        .max()
        .sort_values("decay_score", ascending=False)
        .head(max_page_types_facet)["page_type"]
        .tolist()
    )
    chart_df = top_per_type[top_per_type["page_type"].isin(selected_types)].copy()
    chart_df = chart_df.sort_values(["page_type", "decay_score"], ascending=[True, False])
    chart_df["rank_in_type"] = chart_df.groupby("page_type").cumcount() + 1

    facet_rows = max(1, math.ceil(len(selected_types) / 2))
    fig = px.bar(
        chart_df,
        x="decay_score",
        y="page_short",
        color="severity",
        orientation="h",
        facet_col="page_type",
        facet_col_wrap=2,
        facet_col_spacing=0.08,
        facet_row_spacing=0.08,
        hover_data=["page", "total_revenue", "total_clicks", "losing_revenue_and_clicks", "rank_in_type"],
        category_orders={"severity": SEVERITY_ORDER},
        title=f"Top {top_n_per_type} Risk Pages by Page Type (Top {len(selected_types)} Types)",
        labels={"decay_score": "Decay Score", "page_short": "Page"},
        height=max(950, 430 * facet_rows),
    )
    fig.for_each_annotation(
        lambda annotation: annotation.update(text=annotation.text.replace("page_type=", ""), font={"size": 11})
    )
    fig.update_layout(showlegend=True, legend_title="Severity")
    fig.update_xaxes(matches=None, showticklabels=True)
    fig.update_yaxes(matches=None, showticklabels=True, tickfont={"size": 9})
    apply_base_layout(fig, left_margin=190, width=1600)
    return top_per_type, fig


def build_macro_distribution(scores: pd.DataFrame) -> tuple[go.Figure, go.Figure]:
    hist = px.histogram(
        scores,
        x="decay_score",
        nbins=40,
        title="Risk Score Distribution",
        labels={"decay_score": "Decay Score", "count": "Pages"},
        height=500,
    )
    severity_counts = (
        scores["severity"].value_counts().reindex(SEVERITY_ORDER).fillna(0).reset_index()
    )
    severity_counts.columns = ["severity", "count"]
    sev = px.bar(
        severity_counts,
        x="severity",
        y="count",
        color="severity",
        category_orders={"severity": SEVERITY_ORDER},
        title="Severity Distribution",
        labels={"severity": "Severity", "count": "Pages"},
        height=500,
    )
    apply_base_layout(hist, left_margin=80, height=500, width=1200)
    apply_base_layout(sev, left_margin=80, height=500, width=1200)
    return hist, sev


def build_trend_panel(
    scores: pd.DataFrame,
    monthly: pd.DataFrame,
    *,
    max_pages: int,
) -> tuple[pd.DataFrame, go.Figure]:
    candidates = scores[scores["losing_revenue_and_clicks"]].copy()
    candidates = candidates.sort_values("total_revenue", ascending=False).head(max_pages)
    selected_pages = candidates["page"].tolist()
    if not selected_pages:
        return candidates, go.Figure()

    trends = monthly[monthly["page"].isin(selected_pages)].copy()
    trends["month"] = pd.to_datetime(trends["month"], errors="coerce")
    trends = trends.sort_values(["page", "month"])
    trends["revenue_index"] = trends.groupby("page")["revenue"].transform(normalized_index)
    trends["clicks_index"] = trends.groupby("page")["clicks"].transform(normalized_index)

    rows = len(selected_pages)
    fig = make_subplots(
        rows=rows,
        cols=1,
        shared_xaxes=False,
        subplot_titles=[shorten_page(page, max_len=80) for page in selected_pages],
        vertical_spacing=0.1 if rows > 1 else 0.15,
    )
    for idx, page in enumerate(selected_pages, start=1):
        subset = trends[trends["page"] == page]
        fig.add_trace(
            go.Scatter(
                x=subset["month"],
                y=subset["revenue_index"],
                name="Revenue Index (Base=100)",
                mode="lines+markers",
                line={"color": "#1f77b4"},
                showlegend=(idx == 1),
            ),
            row=idx,
            col=1,
        )
        fig.add_trace(
            go.Scatter(
                x=subset["month"],
                y=subset["clicks_index"],
                name="Clicks Index (Base=100)",
                mode="lines+markers",
                line={"color": "#d62728"},
                showlegend=(idx == 1),
            ),
            row=idx,
            col=1,
        )
        fig.update_yaxes(title_text="Index", row=idx, col=1)
    fig.update_layout(
        title=f"Trend Panels: Most Valuable At-Risk Pages (Top {len(selected_pages)})",
        height=max(500, 320 * rows),
    )
    fig.update_xaxes(title_text="Month")
    apply_base_layout(fig, left_margin=120, width=1300)
    return candidates, fig


def assign_trend_quadrant(scores: pd.DataFrame) -> pd.DataFrame:
    scored = scores.copy()
    scored["slope_clicks"] = pd.to_numeric(scored["slope_clicks"], errors="coerce")
    scored["slope_rev"] = pd.to_numeric(scored["slope_rev"], errors="coerce")
    scored["trend_quadrant"] = np.select(
        [
            (scored["slope_clicks"] < 0) & (scored["slope_rev"] < 0),
            (scored["slope_clicks"] < 0) & (scored["slope_rev"] >= 0),
            (scored["slope_clicks"] >= 0) & (scored["slope_rev"] < 0),
        ],
        [
            "Both Down",
            "Clicks Down, Revenue Stable/Up",
            "Revenue Down, Clicks Stable/Up",
        ],
        default="Both Stable/Up",
    )
    return scored


def build_correlation_outputs(scores: pd.DataFrame) -> tuple[pd.DataFrame, pd.DataFrame, go.Figure, dict[str, float | int | None]]:
    enriched = assign_trend_quadrant(scores)
    valid = enriched[
        pd.to_numeric(enriched["slope_clicks"], errors="coerce").notna()
        & pd.to_numeric(enriched["slope_rev"], errors="coerce").notna()
    ].copy()

    corr_pearson: float | None = None
    corr_spearman: float | None = None
    if len(valid) >= 3:
        corr_pearson = float(valid["slope_clicks"].corr(valid["slope_rev"], method="pearson"))
        # Avoid scipy dependency by computing Spearman as Pearson on rank-transformed values.
        corr_spearman = float(
            valid["slope_clicks"].rank(method="average").corr(
                valid["slope_rev"].rank(method="average"),
                method="pearson",
            )
        )

    quadrant_summary = (
        valid.groupby("trend_quadrant", as_index=False)
        .agg(
            pages=("page", "count"),
            total_revenue=("total_revenue", "sum"),
            total_clicks=("total_clicks", "sum"),
            avg_decay_score=("decay_score", "mean"),
        )
    )
    quadrant_summary["trend_quadrant"] = pd.Categorical(
        quadrant_summary["trend_quadrant"],
        categories=QUADRANT_ORDER,
        ordered=True,
    )
    quadrant_summary = quadrant_summary.sort_values("trend_quadrant").reset_index(drop=True)
    total_pages = int(quadrant_summary["pages"].sum()) if not quadrant_summary.empty else 0
    quadrant_summary["share_pages_pct"] = np.where(
        total_pages > 0,
        (quadrant_summary["pages"] / total_pages) * 100.0,
        0.0,
    )

    fig = px.scatter(
        valid,
        x="slope_clicks",
        y="slope_rev",
        color="trend_quadrant",
        size="total_revenue",
        size_max=24,
        hover_data=["page", "page_type", "decay_score", "total_revenue", "total_clicks", "severity"],
        category_orders={"trend_quadrant": QUADRANT_ORDER},
        title="Clicks Trend vs Revenue Trend (Page Level)",
        labels={"slope_clicks": "Clicks Slope", "slope_rev": "Revenue Slope", "trend_quadrant": "Quadrant"},
        height=640,
    )
    fig.add_hline(y=0.0, line_width=1, line_dash="dash", line_color="#666666")
    fig.add_vline(x=0.0, line_width=1, line_dash="dash", line_color="#666666")
    apply_base_layout(fig, left_margin=90, height=640, width=1300)

    metrics: dict[str, float | int | None] = {
        "rows_with_valid_slopes": int(len(valid)),
        "pearson_slope_clicks_vs_rev": corr_pearson,
        "spearman_slope_clicks_vs_rev": corr_spearman,
    }
    return valid, quadrant_summary, fig, metrics


def build_quadrant_page_exports(valid_scores: pd.DataFrame) -> dict[str, pd.DataFrame]:
    ordered = valid_scores.sort_values("decay_score", ascending=False).copy()
    cols = [
        "page",
        "page_type",
        "decay_score",
        "severity",
        "slope_clicks",
        "slope_rev",
        "total_revenue",
        "total_clicks",
        "trend_quadrant",
    ]
    for col in cols:
        if col not in ordered.columns:
            ordered[col] = pd.NA
    export_cols = ordered[cols].copy()
    return {
        "pages_both_down": export_cols[export_cols["trend_quadrant"] == "Both Down"].copy(),
        "pages_clicks_down_revenue_stable_up": export_cols[
            export_cols["trend_quadrant"] == "Clicks Down, Revenue Stable/Up"
        ].copy(),
        "pages_revenue_down_clicks_stable_up": export_cols[
            export_cols["trend_quadrant"] == "Revenue Down, Clicks Stable/Up"
        ].copy(),
    }


def build_lumar_page_diagnostics(
    scores: pd.DataFrame,
    *,
    lumar_pages_csv: Path,
    top_focus_n: int,
) -> tuple[pd.DataFrame, pd.DataFrame, go.Figure, dict[str, float | int | str | None]]:
    df = pd.read_csv(lumar_pages_csv)
    if df.empty:
        empty = pd.DataFrame()
        return empty, empty, go.Figure(), {"rows_lumar_input": 0}

    url_col = pick_first_column(list(df.columns), ["url", "page", "page_url", "address"])
    if not url_col:
        raise ValueError(
            f"{lumar_pages_csv} must include a page URL column (url/page/page_url/address)."
        )

    work = df.copy()
    work["page"] = work[url_col].map(normalize_page_path)
    work = work[work["page"] != ""].copy()

    numeric_candidates = {
        "deeprank": ["deeprank", "deep_rank"],
        "level": ["level", "crawl_depth", "depth"],
        "unique_links_in_count": ["unique_links_in_count", "unique_inlinks", "inlinks_unique"],
        "links_in_count": ["links_in_count", "inlinks", "internal_links_in"],
        "word_count": ["word_count", "words"],
    }
    for out_col, candidates in numeric_candidates.items():
        source = pick_first_column(list(work.columns), candidates)
        if source:
            work[out_col] = pd.to_numeric(work[source], errors="coerce")
        else:
            work[out_col] = np.nan

    sitemap_col = pick_first_column(list(work.columns), ["found_in_sitemap", "in_sitemap", "is_in_sitemap"])
    crawl_col = pick_first_column(list(work.columns), ["found_in_web_crawl", "in_crawl", "is_in_crawl"])
    if sitemap_col:
        work["found_in_sitemap"] = to_bool_series(work[sitemap_col])
    else:
        work["found_in_sitemap"] = False
    if crawl_col:
        work["found_in_web_crawl"] = to_bool_series(work[crawl_col])
    else:
        work["found_in_web_crawl"] = False

    lumar_page = (
        work.groupby("page", as_index=False)
        .agg(
            {
                "deeprank": "max",
                "level": "min",
                "unique_links_in_count": "max",
                "links_in_count": "max",
                "word_count": "max",
                "found_in_sitemap": "max",
                "found_in_web_crawl": "max",
            }
        )
    )

    low_links_threshold = float(lumar_page["unique_links_in_count"].quantile(0.25))
    low_deeprank_threshold = float(lumar_page["deeprank"].quantile(0.25))
    deep_level_threshold = int(max(4, math.ceil(float(lumar_page["level"].quantile(0.75)))))

    merged = scores.merge(lumar_page, on="page", how="left")
    merged["found_in_sitemap"] = to_bool_series(merged["found_in_sitemap"])
    merged["found_in_web_crawl"] = to_bool_series(merged["found_in_web_crawl"])
    merged["lumar_has_data"] = merged["deeprank"].notna() | merged["level"].notna() | merged["unique_links_in_count"].notna()
    merged["lumar_issue_not_in_sitemap"] = merged["lumar_has_data"] & (~merged["found_in_sitemap"])
    merged["lumar_issue_low_internal_links"] = (
        merged["lumar_has_data"]
        & merged["unique_links_in_count"].notna()
        & (merged["unique_links_in_count"] <= low_links_threshold)
    )
    merged["lumar_issue_deep_level"] = (
        merged["lumar_has_data"] & merged["level"].notna() & (merged["level"] >= deep_level_threshold)
    )
    merged["lumar_issue_low_deeprank"] = (
        merged["lumar_has_data"] & merged["deeprank"].notna() & (merged["deeprank"] <= low_deeprank_threshold)
    )
    merged["lumar_issue_count"] = (
        merged[
            [
                "lumar_issue_not_in_sitemap",
                "lumar_issue_low_internal_links",
                "lumar_issue_deep_level",
                "lumar_issue_low_deeprank",
            ]
        ]
        .astype(int)
        .sum(axis=1)
    )
    merged["lumar_has_issue"] = merged["lumar_issue_count"] > 0

    focus = merged.sort_values("decay_score", ascending=False).head(top_focus_n).copy()
    issue_summary = pd.DataFrame(
        {
            "issue_category": [
                "Not in sitemap",
                "Low internal links",
                "Deep crawl level",
                "Low deeprank",
            ],
            "count": [
                int(focus["lumar_issue_not_in_sitemap"].sum()),
                int(focus["lumar_issue_low_internal_links"].sum()),
                int(focus["lumar_issue_deep_level"].sum()),
                int(focus["lumar_issue_low_deeprank"].sum()),
            ],
        }
    ).sort_values("count", ascending=False)

    fig = px.bar(
        issue_summary.sort_values("count", ascending=True),
        x="count",
        y="issue_category",
        orientation="h",
        title=f"Lumar Diagnostic Signals Across Top {top_focus_n} Risk Pages",
        labels={"count": "Pages", "issue_category": "Lumar Diagnostic Signal"},
        height=520,
    )
    apply_base_layout(fig, left_margin=210, height=520, width=1200)

    meta: dict[str, float | int | str | None] = {
        "rows_lumar_input": int(len(df)),
        "lumar_unique_pages": int(lumar_page["page"].nunique()),
        "scored_pages_with_lumar_match": int(merged["lumar_has_data"].sum()),
        "low_links_threshold_p25": low_links_threshold,
        "low_deeprank_threshold_p25": low_deeprank_threshold,
        "deep_level_threshold": deep_level_threshold,
    }
    return merged, issue_summary, fig, meta


def build_lumar_issue_chart(lumar_csv: Path) -> tuple[pd.DataFrame, go.Figure]:
    df = pd.read_csv(lumar_csv)
    if df.empty:
        return pd.DataFrame(columns=["issue_category", "count"]), go.Figure()

    columns = list(df.columns)
    lower_map = {column.lower(): column for column in columns}
    category_col = None
    for candidate in ["issue_category", "category", "issue", "issue type", "issue_type", "name"]:
        if candidate.lower() in lower_map:
            category_col = lower_map[candidate.lower()]
            break
    if not category_col:
        raise ValueError(
            f"{lumar_csv} needs an issue category column "
            "(e.g. issue_category/category/issue)."
        )

    count_col = None
    for candidate in ["count", "urls", "url_count", "affected_urls", "pages"]:
        if candidate.lower() in lower_map:
            count_col = lower_map[candidate.lower()]
            break

    grouped = df.copy()
    if count_col:
        grouped["count"] = pd.to_numeric(grouped[count_col], errors="coerce").fillna(0.0)
        grouped = grouped.groupby(category_col, as_index=False)["count"].sum()
    else:
        grouped = grouped.groupby(category_col, as_index=False).size().rename(columns={"size": "count"})

    grouped = grouped.rename(columns={category_col: "issue_category"})
    grouped = grouped.sort_values("count", ascending=False).head(20)

    fig = px.bar(
        grouped.sort_values("count", ascending=True),
        x="count",
        y="issue_category",
        orientation="h",
        title="Lumar Issue Categories (Top 20)",
        labels={"count": "Count", "issue_category": "Issue Category"},
        height=max(500, 24 * len(grouped)),
    )
    apply_base_layout(fig, left_margin=220, width=1200)
    return grouped, fig


def _safe_float(value: object) -> float | None:
    try:
        if value is None:
            return None
        if isinstance(value, str) and value.strip() == "":
            return None
        parsed = float(value)
        if math.isnan(parsed):
            return None
        return parsed
    except Exception:
        return None


def build_pptx_report(
    *,
    output_pptx: Path,
    title_text: str,
    summary: dict[str, object],
    generated_chart_files: dict[str, str],
    top_overall_df: pd.DataFrame,
    quadrant_summary_df: pd.DataFrame,
    correlation_metrics: dict[str, float | int | None],
    lumar_focus_df: pd.DataFrame | None,
    lumar_issue_summary_df: pd.DataFrame | None,
) -> bool:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:  # pragma: no cover - environment dependent
        print(f"Warning: python-pptx unavailable, skipping PPTX export: {exc}")
        return False

    def add_picture_fit(
        slide,
        image_path: Path,
        *,
        box_left,
        box_top,
        box_width,
        box_height,
    ) -> None:
        """
        Add an image to a slide, scaled to fit a bounding box while preserving aspect ratio.
        The image is centered in the box and never allowed to overflow it.
        """
        picture = slide.shapes.add_picture(str(image_path), box_left, box_top)
        original_width = int(picture.width)
        original_height = int(picture.height)
        if original_width <= 0 or original_height <= 0:
            return

        max_width = int(box_width)
        max_height = int(box_height)
        width_scale = max_width / original_width
        height_scale = max_height / original_height
        scale = min(width_scale, height_scale)

        new_width = int(original_width * scale)
        new_height = int(original_height * scale)
        picture.width = new_width
        picture.height = new_height
        picture.left = int(box_left + (max_width - new_width) / 2)
        picture.top = int(box_top + (max_height - new_height) / 2)

    def get_content_box(
        slide,
        *,
        side_margin_in: float = 0.4,
        title_gap_in: float = 0.15,
        top_fallback_in: float = 1.0,
        bottom_margin_in: float = 0.3,
    ) -> tuple[int, int, int, int]:
        """
        Return a safe content box that stays below the title placeholder and inside slide bounds.
        """
        slide_width = int(prs.slide_width)
        slide_height = int(prs.slide_height)
        left = int(Inches(side_margin_in))
        right_margin = int(Inches(side_margin_in))
        title_bottom = 0
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None:
            title_bottom = int(title_shape.top + title_shape.height)
        top = max(int(Inches(top_fallback_in)), title_bottom + int(Inches(title_gap_in)))
        bottom_margin = int(Inches(bottom_margin_in))
        width = max(1, slide_width - left - right_margin)
        height = max(1, slide_height - top - bottom_margin)
        return left, top, width, height

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title_text
    subtitle_text = (
        f"Run: {summary.get('run_date_utc', 'n/a')}\n"
        f"Pages scored: {summary.get('counts', {}).get('pages_scored', 'n/a')}\n"
        f"Pages losing revenue and clicks: "
        f"{summary.get('counts', {}).get('pages_losing_revenue_and_clicks', 'n/a')}"
    )
    title_slide.placeholders[1].text = subtitle_text

    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "Attention Summary"
    text_frame = summary_slide.placeholders[1].text_frame
    text_frame.clear()
    lines = [
        f"Analysis window: {summary.get('analysis_window', {}).get('min_month')} to {summary.get('analysis_window', {}).get('max_month')}",
        f"Top at-risk pages exported: {len(top_overall_df)}",
        f"Slope correlation (Pearson): {(_safe_float(correlation_metrics.get('pearson_slope_clicks_vs_rev')) or 0):.3f}",
        f"Slope correlation (Spearman): {(_safe_float(correlation_metrics.get('spearman_slope_clicks_vs_rev')) or 0):.3f}",
    ]
    for idx, line in enumerate(lines):
        if idx == 0:
            text_frame.text = line
            text_frame.paragraphs[0].font.size = Pt(18)
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.font.size = Pt(16)

    table_slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[6])
    if table_slide.shapes.title:
        table_slide.shapes.title.text = "Top Pages Requiring Attention"
    table_source = top_overall_df.sort_values("decay_score", ascending=False).head(12).copy()
    table_source["page"] = table_source["page"].map(lambda p: shorten_page(p, max_len=70))
    table_columns = ["page", "page_type", "decay_score", "severity", "slope_rev", "slope_clicks"]
    table_source = table_source[table_columns].copy()
    rows = len(table_source) + 1
    cols = len(table_columns)
    table_left, table_top, table_width, table_height = get_content_box(table_slide)
    table_shape = table_slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = table_shape.table
    col_width_ratios = [0.46, 0.16, 0.11, 0.10, 0.085, 0.085]
    for idx, ratio in enumerate(col_width_ratios):
        table.columns[idx].width = int(table_width * ratio)
    for col_idx, header in enumerate(table_columns):
        table.cell(0, col_idx).text = header
    for row_idx, (_, row) in enumerate(table_source.iterrows(), start=1):
        for col_idx, header in enumerate(table_columns):
            value = row[header]
            if isinstance(value, float):
                if header == "decay_score":
                    text = f"{value:.2f}"
                else:
                    text = f"{value:.4f}"
            else:
                text = str(value)
            table.cell(row_idx, col_idx).text = text
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            if not cell.text_frame.paragraphs:
                continue
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10 if row_idx > 0 else 11)

    image_sequence = [
        ("top_20_risk_pages_overall_png", "Top Risk Pages"),
        ("top_10_risk_pages_by_page_type_png", "Top Risk Pages by Page Type"),
        ("risk_score_distribution_png", "Risk Score Distribution"),
        ("severity_distribution_png", "Severity Distribution"),
        ("trend_panels_png", "Most Valuable At-Risk Trend Panels"),
        ("trend_correlation_scatter_png", "Click vs Revenue Trend Correlation"),
        ("quadrant_summary_png", "Trend Quadrant Distribution"),
    ]
    for key, slide_title in image_sequence:
        image_path = generated_chart_files.get(key)
        if not image_path:
            continue
        image_file = Path(image_path)
        if not image_file.exists():
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[6])
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
        box_left, box_top, box_width, box_height = get_content_box(slide)
        add_picture_fit(
            slide,
            image_file,
            box_left=box_left,
            box_top=box_top,
            box_width=box_width,
            box_height=box_height,
        )

    quadrant_slide = prs.slides.add_slide(prs.slide_layouts[1])
    quadrant_slide.shapes.title.text = "Clicks vs Revenue Quadrants"
    q_text = quadrant_slide.placeholders[1].text_frame
    q_text.clear()
    if quadrant_summary_df.empty:
        q_text.text = "No valid slope rows available."
    else:
        for idx, (_, row) in enumerate(quadrant_summary_df.iterrows()):
            line = (
                f"{row['trend_quadrant']}: {int(row['pages'])} pages "
                f"({float(row['share_pages_pct']):.1f}% of valid pages)"
            )
            if idx == 0:
                q_text.text = line
                q_text.paragraphs[0].font.size = Pt(15)
            else:
                para = q_text.add_paragraph()
                para.text = line
                para.font.size = Pt(15)

    # Lumar slides intentionally excluded for now.

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return True


def run(args: argparse.Namespace) -> None:
    scores_csv = Path(args.scores_csv)
    monthly_csv = Path(args.monthly_adjusted_csv)
    if not scores_csv.exists():
        raise FileNotFoundError(f"Scores CSV not found: {scores_csv}")
    if not monthly_csv.exists():
        raise FileNotFoundError(f"Monthly adjusted CSV not found: {monthly_csv}")

    scores = pd.read_csv(scores_csv)
    monthly = pd.read_csv(monthly_csv)
    require_columns(
        scores,
        [
            "page",
            "page_type",
            "decay_score",
            "severity",
            "total_revenue",
            "total_clicks",
            "losing_revenue_and_clicks",
            "slope_rev",
            "slope_clicks",
        ],
        file_label=str(scores_csv),
    )
    require_columns(
        monthly,
        ["page", "month", "revenue", "clicks"],
        file_label=str(monthly_csv),
    )
    monthly["month"] = pd.to_datetime(monthly["month"], errors="coerce")

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    charts_dir = output_dir / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)
    exports_dir = output_dir / "exports"
    exports_dir.mkdir(parents=True, exist_ok=True)

    generated: dict[str, str] = {}
    top_overall_label = f"top_{args.top_n_overall}_risk_pages_overall"
    top_by_type_label = f"top_{args.top_n_per_page_type}_risk_pages_by_page_type"

    top_overall_df, top_overall_fig = build_top_overall(scores, top_n=args.top_n_overall)
    top_overall_csv = exports_dir / f"{top_overall_label}.csv"
    top_overall_png = charts_dir / f"{top_overall_label}.png"
    top_overall_df.to_csv(top_overall_csv, index=False)
    if write_plot_image(top_overall_fig, top_overall_png):
        generated["top_20_risk_pages_overall_png"] = str(top_overall_png)

    top_per_type_df, top_per_type_fig = build_top_per_page_type(
        scores,
        top_n_per_type=args.top_n_per_page_type,
        max_page_types_facet=args.max_page_types_facet,
    )
    top_per_type_csv = exports_dir / f"{top_by_type_label}.csv"
    top_per_type_png = charts_dir / f"{top_by_type_label}.png"
    top_per_type_df.to_csv(top_per_type_csv, index=False)
    if not top_per_type_df.empty and write_plot_image(top_per_type_fig, top_per_type_png):
        generated["top_10_risk_pages_by_page_type_png"] = str(top_per_type_png)

    hist_fig, severity_fig = build_macro_distribution(scores)
    risk_dist_png = charts_dir / "risk_score_distribution.png"
    severity_png = charts_dir / "severity_distribution.png"
    if write_plot_image(hist_fig, risk_dist_png):
        generated["risk_score_distribution_png"] = str(risk_dist_png)
    if write_plot_image(severity_fig, severity_png):
        generated["severity_distribution_png"] = str(severity_png)

    trend_pages_df, trend_fig = build_trend_panel(
        scores,
        monthly,
        max_pages=args.trend_pages_count,
    )
    trend_csv = exports_dir / "most_valuable_at_risk_pages.csv"
    trend_png = charts_dir / "trend_panels_most_valuable_at_risk.png"
    trend_pages_df.to_csv(trend_csv, index=False)
    if not trend_pages_df.empty and write_plot_image(trend_fig, trend_png):
        generated["trend_panels_png"] = str(trend_png)

    valid_scores, quadrant_summary, corr_fig, correlation_metrics = build_correlation_outputs(scores)
    corr_png = charts_dir / "trend_correlation_scatter.png"
    quadrant_summary_csv = exports_dir / "trend_quadrant_summary.csv"
    quadrant_summary.to_csv(quadrant_summary_csv, index=False)
    if write_plot_image(corr_fig, corr_png):
        generated["trend_correlation_scatter_png"] = str(corr_png)

    quadrant_bar_fig = px.bar(
        quadrant_summary,
        x="trend_quadrant",
        y="pages",
        color="trend_quadrant",
        category_orders={"trend_quadrant": QUADRANT_ORDER},
        title="Page Count by Click/Revenue Trend Quadrant",
        labels={"trend_quadrant": "Quadrant", "pages": "Pages"},
        height=520,
    )
    apply_base_layout(quadrant_bar_fig, left_margin=90, height=520, width=1200)
    quadrant_bar_png = charts_dir / "trend_quadrant_distribution.png"
    if write_plot_image(quadrant_bar_fig, quadrant_bar_png):
        generated["quadrant_summary_png"] = str(quadrant_bar_png)

    quadrant_exports = build_quadrant_page_exports(valid_scores)
    quadrant_csv_map = {
        "pages_both_down": exports_dir / "pages_both_down.csv",
        "pages_clicks_down_revenue_stable_up": exports_dir / "pages_clicks_down_revenue_stable_up.csv",
        "pages_revenue_down_clicks_stable_up": exports_dir / "pages_revenue_down_clicks_stable_up.csv",
    }
    for key, path in quadrant_csv_map.items():
        quadrant_exports[key].to_csv(path, index=False)

    lumar_summary: dict[str, object] = {
        "issue_categories": {"used": False},
        "page_diagnostics": {"used": False},
    }
    lumar_focus_for_pptx: pd.DataFrame | None = None
    lumar_issue_summary_for_pptx: pd.DataFrame | None = None

    if args.lumar_pages_csv:
        lumar_pages_csv = Path(args.lumar_pages_csv)
        if not lumar_pages_csv.exists():
            raise FileNotFoundError(f"Lumar pages CSV not found: {lumar_pages_csv}")
        lumar_joined, lumar_issue_summary, lumar_issue_fig, lumar_meta = build_lumar_page_diagnostics(
            scores,
            lumar_pages_csv=lumar_pages_csv,
            top_focus_n=args.top_n_overall,
        )
        lumar_focus = lumar_joined.sort_values("decay_score", ascending=False).head(args.top_n_overall).copy()
        lumar_focus_csv = exports_dir / f"{top_overall_label}_with_lumar_diagnostics.csv"
        lumar_issue_summary_csv = exports_dir / "lumar_diagnostic_issue_summary.csv"
        lumar_issue_signals_png = charts_dir / "lumar_diagnostic_issue_signals.png"
        lumar_focus.to_csv(lumar_focus_csv, index=False)
        lumar_issue_summary.to_csv(lumar_issue_summary_csv, index=False)
        if write_plot_image(lumar_issue_fig, lumar_issue_signals_png):
            generated["lumar_issue_signals_png"] = str(lumar_issue_signals_png)

        lumar_summary["page_diagnostics"] = {
            "used": True,
            "input_csv": str(lumar_pages_csv),
            "output_focus_csv": str(lumar_focus_csv),
            "output_issue_summary_csv": str(lumar_issue_summary_csv),
            "output_issue_signals_png": str(lumar_issue_signals_png),
            **lumar_meta,
        }
        lumar_focus_for_pptx = lumar_focus[lumar_focus["lumar_has_issue"]].copy()
        lumar_issue_summary_for_pptx = lumar_issue_summary.copy()

    if args.lumar_issues_csv:
        lumar_csv = Path(args.lumar_issues_csv)
        if not lumar_csv.exists():
            raise FileNotFoundError(f"Lumar issues CSV not found: {lumar_csv}")
        lumar_df, lumar_fig = build_lumar_issue_chart(lumar_csv)
        lumar_csv_out = exports_dir / "lumar_issue_categories_top20.csv"
        lumar_png = charts_dir / "lumar_issue_categories.png"
        lumar_df.to_csv(lumar_csv_out, index=False)
        lumar_summary["issue_categories"] = {
            "used": True,
            "input_csv": str(lumar_csv),
            "rows": int(len(lumar_df)),
            "output_csv": str(lumar_csv_out),
            "output_png": str(lumar_png),
        }
        if not lumar_df.empty and write_plot_image(lumar_fig, lumar_png):
            generated["lumar_issue_categories_png"] = str(lumar_png)

    summary = {
        "run_date_utc": datetime.utcnow().isoformat() + "Z",
        "inputs": {
            "scores_csv": str(scores_csv),
            "monthly_adjusted_csv": str(monthly_csv),
            "lumar_issues_csv": args.lumar_issues_csv or None,
            "lumar_pages_csv": args.lumar_pages_csv or None,
        },
        "settings": {
            "top_n_overall": args.top_n_overall,
            "top_n_per_page_type": args.top_n_per_page_type,
            "max_page_types_facet": args.max_page_types_facet,
            "trend_pages_count": args.trend_pages_count,
            "export_pptx": bool(args.export_pptx),
        },
        "analysis_window": {
            "min_month": monthly["month"].min().date().isoformat() if not monthly.empty else None,
            "max_month": monthly["month"].max().date().isoformat() if not monthly.empty else None,
            "month_count": int(monthly["month"].nunique()) if not monthly.empty else 0,
        },
        "counts": {
            "pages_scored": int(len(scores)),
            "page_types": int(scores["page_type"].nunique()),
            "pages_losing_revenue_and_clicks": int(
                pd.to_numeric(scores["losing_revenue_and_clicks"], errors="coerce").fillna(0).astype(bool).sum()
            ),
            "trend_pages_selected": int(len(trend_pages_df)),
        },
        "correlation": {
            **correlation_metrics,
            "quadrant_counts": (
                quadrant_summary.set_index("trend_quadrant")["pages"].to_dict()
                if not quadrant_summary.empty
                else {}
            ),
        },
        "outputs": {
            "charts_dir": str(charts_dir),
            "exports_dir": str(exports_dir),
            "top_20_csv": str(top_overall_csv),
            "top_10_by_type_csv": str(top_per_type_csv),
            "most_valuable_at_risk_csv": str(trend_csv),
            "trend_quadrant_summary_csv": str(quadrant_summary_csv),
            "quadrant_export_csvs": {k: str(v) for k, v in quadrant_csv_map.items()},
            "generated_chart_files": generated,
        },
        "lumar": lumar_summary,
    }

    if args.export_pptx:
        pptx_output = Path(args.pptx_output) if args.pptx_output else output_dir / "revenue_decay_attention_pack.pptx"
        wrote_pptx = build_pptx_report(
            output_pptx=pptx_output,
            title_text=args.pptx_title,
            summary=summary,
            generated_chart_files=generated,
            top_overall_df=top_overall_df,
            quadrant_summary_df=quadrant_summary,
            correlation_metrics=correlation_metrics,
            lumar_focus_df=lumar_focus_for_pptx,
            lumar_issue_summary_df=lumar_issue_summary_for_pptx,
        )
        if wrote_pptx:
            summary["outputs"]["pptx"] = str(pptx_output)

    summary_out = output_dir / "summary.json"
    summary_out.write_text(json.dumps(summary, indent=2), encoding="utf-8")

    print(f"Wrote: {top_overall_csv}")
    print(f"Wrote: {top_per_type_csv}")
    print(f"Wrote: {trend_csv}")
    print(f"Wrote: {quadrant_summary_csv}")
    print(f"Wrote: {summary_out}")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Generate a focused chart pack from revenue decay scoring outputs."
    )
    parser.add_argument(
        "--scores-csv",
        default=str(
            ROOT_DIR
            / "reports"
            / "analysis"
            / "revenue_decay_radar"
            / "scoring"
            / "decay_scores_all_pages.csv"
        ),
        help="Scored pages CSV from revenue_decay_scoring.",
    )
    parser.add_argument(
        "--monthly-adjusted-csv",
        default=str(
            ROOT_DIR
            / "reports"
            / "analysis"
            / "revenue_decay_radar"
            / "scoring"
            / "monthly_adjusted_signals.csv"
        ),
        help="Monthly adjusted signals CSV from revenue_decay_scoring.",
    )
    parser.add_argument(
        "--lumar-issues-csv",
        default="",
        help="Optional Lumar issue-category CSV for issue chart.",
    )
    parser.add_argument(
        "--lumar-pages-csv",
        default="",
        help=(
            "Optional Lumar page export CSV (e.g., lumar_bfb.csv) for page-level "
            "diagnostic join with scored pages."
        ),
    )
    parser.add_argument(
        "--top-n-overall",
        type=int,
        default=20,
        help="Top N risk pages overall.",
    )
    parser.add_argument(
        "--top-n-per-page-type",
        type=int,
        default=10,
        help="Top N risk pages per page_type.",
    )
    parser.add_argument(
        "--max-page-types-facet",
        type=int,
        default=8,
        help="How many page types to include in per-type facet chart.",
    )
    parser.add_argument(
        "--trend-pages-count",
        type=int,
        default=2,
        help="How many most-valuable at-risk pages to include in trend panel.",
    )
    parser.add_argument(
        "--output-dir",
        default=str(
            ROOT_DIR
            / "reports"
            / "analysis"
            / "revenue_decay_radar"
            / "chart_pack"
        ),
        help="Output directory for charts and exports.",
    )
    parser.add_argument(
        "--export-pptx",
        action="store_true",
        help="If set, export a PPTX attention pack with charts and key summaries.",
    )
    parser.add_argument(
        "--pptx-output",
        default="",
        help="Optional output PPTX path. Defaults to <output-dir>/revenue_decay_attention_pack.pptx.",
    )
    parser.add_argument(
        "--pptx-title",
        default="Revenue Decay Radar: Pages Requiring Attention",
        help="Title text for the generated PPTX.",
    )
    return parser


if __name__ == "__main__":
    parser = build_parser()
    run(parser.parse_args())
